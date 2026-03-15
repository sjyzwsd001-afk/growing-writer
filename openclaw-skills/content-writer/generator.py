#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
内容写作助手 - Word & PPT 生成
"""

import sys
import os
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
import argparse
from datetime import datetime

def create_word_doc(title, outline, content, output_path):
    """生成 Word 文档"""
    doc = Document()
    
    # 标题
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 日期
    date_para = doc.add_paragraph(datetime.now().strftime('%Y年%m月%d日'))
    date_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    
    doc.add_paragraph()
    
    # 大纲
    doc.add_heading('大纲', level=1)
    for i, item in enumerate(outline, 1):
        doc.add_paragraph(f"{i}. {item}", style='List Number')
    
    doc.add_paragraph()
    
    # 正文内容
    doc.add_heading('正文', level=1)
    for section_title, section_content in content:
        doc.add_heading(section_title, level=2)
        doc.add_paragraph(section_content)
    
    # 保存
    doc.save(output_path)
    return output_path

def create_ppt(title, outline, content, output_path):
    """生成 PPT"""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    # 封面页
    slide_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12), PptxInches(1.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = PptxPt(44)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(4.5), PptxInches(12), PptxInches(0.5))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = datetime.now().strftime('%Y年%m月%d日')
    date_p.font.size = PptxPt(18)
    date_p.alignment = PP_ALIGN.CENTER
    
    # 大纲页
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "目录"
    title_p.font.size = PptxPt(32)
    title_p.font.bold = True
    
    content_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.5), PptxInches(11), PptxInches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(outline, 1):
        p = content_frame.paragraphs[0] if i == 1 else content_frame.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.size = PptxPt(20)
        p.space_after = PptxPt(12)
    
    # 内容页
    for section_title, section_content in content:
        slide = prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = section_title
        title_p.font.size = PptxPt(28)
        title_p.font.bold = True
        
        # 内容
        content_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.5), PptxInches(11), PptxInches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        p = content_frame.paragraphs[0]
        p.text = section_content[:500]  # 限制长度
        p.font.size = PptxPt(16)
        p.space_after = PptxPt(12)
    
    # 保存
    prs.save(output_path)
    return output_path

def main():
    parser = argparse.ArgumentParser(description='内容写作助手')
    parser.add_argument('--mode', choices=['word', 'ppt'], required=True, help='输出模式')
    parser.add_argument('--title', required=True, help='文档标题')
    parser.add_argument('--outline', nargs='+', required=True, help='大纲项目')
    parser.add_argument('--content', nargs='+', required=True, help='内容（格式：标题1:内容1 标题2:内容2）')
    parser.add_argument('--output', required=True, help='输出文件路径')
    
    args = parser.parse_args()
    
    # 解析内容
    content = []
    for item in args.content:
        if ':' in item:
            idx = item.index(':')
            section_title = item[:idx].strip()
            section_content = item[idx+1:].strip()
            content.append((section_title, section_content))
    
    if args.mode == 'word':
        create_word_doc(args.title, args.outline, content, args.output)
    else:
        create_ppt(args.title, args.outline, content, args.output)
    
    print(f"生成成功: {args.output}")

if __name__ == '__main__':
    main()
